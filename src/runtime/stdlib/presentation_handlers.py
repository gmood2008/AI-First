from __future__ import annotations

import hashlib
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple

from pptx import Presentation

from ..handler import ActionHandler


def _sha256_file(path: Path) -> str:
    h = hashlib.sha256()
    with path.open("rb") as f:
        for chunk in iter(lambda: f.read(1024 * 1024), b""):
            h.update(chunk)
    return h.hexdigest()


def _contains_external_url(value: Any) -> bool:
    if isinstance(value, str):
        s = value.lower()
        return "http://" in s or "https://" in s
    if isinstance(value, dict):
        return any(_contains_external_url(v) for v in value.values())
    if isinstance(value, list):
        return any(_contains_external_url(v) for v in value)
    return False


def _require_keys(obj: Dict[str, Any], required: List[str], path: str) -> None:
    for k in required:
        if k not in obj:
            raise ValueError(f"Missing required field: {path}.{k}")


def _only_allow_keys(obj: Dict[str, Any], allowed: List[str], path: str) -> None:
    for k in obj.keys():
        if k not in allowed:
            raise ValueError(f"Forbidden field: {path}.{k}")


def _resolve_output_path(context: Any, raw_path: str) -> Path:
    if not raw_path or not isinstance(raw_path, str):
        raise ValueError("output.path is required")

    p = Path(raw_path)
    if not p.is_absolute():
        if not hasattr(context, "workspace_root"):
            raise ValueError("workspace_root is required in execution context")
        p = Path(context.workspace_root) / p

    resolved = p.expanduser().resolve()

    workspace_root = Path(context.workspace_root).resolve()
    try:
        resolved.relative_to(workspace_root)
    except ValueError as e:
        raise ValueError("output.path must be within workspace_root") from e

    if resolved.suffix.lower() != ".pptx":
        raise ValueError("output.path must end with .pptx")

    return resolved


class PPTXRenderHandler(ActionHandler):
    def execute(self, params: Dict[str, Any], context: Any) -> Dict[str, Any]:
        self.validate_params(params)

        allowed_top = {"document", "output", "theme"}
        _only_allow_keys(params, sorted(list(allowed_top)), path="input")

        document = params.get("document")
        output = params.get("output")
        if not isinstance(document, dict):
            raise ValueError("document must be an object")
        if not isinstance(output, dict):
            raise ValueError("output must be an object")

        _require_keys(document, ["title", "sections"], path="document")
        _only_allow_keys(document, ["title", "subtitle", "meta", "sections"], path="document")

        if _contains_external_url(document):
            raise ValueError("input_constraints.forbid_external_urls violated")

        title = document.get("title")
        subtitle = document.get("subtitle")
        if not isinstance(title, str) or not title.strip():
            raise ValueError("document.title must be a non-empty string")
        if subtitle is not None and not isinstance(subtitle, str):
            raise ValueError("document.subtitle must be a string")

        sections = document.get("sections")
        if not isinstance(sections, list) or len(sections) < 1:
            raise ValueError("document.sections must be a non-empty array")

        max_slide_count = int(self.contracts.get("max_slide_count", 50))
        max_file_size_mb = int(self.contracts.get("max_file_size_mb", 20))

        out_path = _resolve_output_path(context, str((output.get("path") or "")).strip())
        out_path.parent.mkdir(parents=True, exist_ok=True)

        prs = Presentation()

        # Title slide
        title_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_layout)
        slide.shapes.title.text = title
        if subtitle and len(slide.placeholders) > 1:
            slide.placeholders[1].text = subtitle

        render_log: List[str] = []
        render_log.append("render:title_slide")

        # Section slides
        for idx, section in enumerate(sections):
            if not isinstance(section, dict):
                raise ValueError(f"document.sections[{idx}] must be an object")
            _require_keys(section, ["heading", "content"], path=f"document.sections[{idx}]")
            _only_allow_keys(section, ["heading", "content"], path=f"document.sections[{idx}]")

            heading = section.get("heading")
            content = section.get("content")
            if not isinstance(heading, str) or not heading.strip():
                raise ValueError(f"document.sections[{idx}].heading must be a non-empty string")
            if not isinstance(content, list) or len(content) < 1:
                raise ValueError(f"document.sections[{idx}].content must be a non-empty array")

            if len(prs.slides) >= max_slide_count:
                raise ValueError("execution_constraints.max_slide_count exceeded")

            layout = prs.slide_layouts[1]
            s = prs.slides.add_slide(layout)
            s.shapes.title.text = heading

            body = s.shapes.placeholders[1].text_frame
            body.clear()

            for j, block in enumerate(content):
                if not isinstance(block, dict):
                    raise ValueError(f"document.sections[{idx}].content[{j}] must be an object")
                if "type" not in block:
                    raise ValueError(f"document.sections[{idx}].content[{j}].type is required")

                block_type = block.get("type")
                if block_type not in {"bullet_list", "paragraph", "table", "citation", "image", "chart"}:
                    raise ValueError(f"Unsupported content type: {block_type}")

                if block_type in {"image", "chart"}:
                    raise ValueError(f"Content type '{block_type}' is forbidden in this renderer")

                if block_type == "paragraph":
                    text = block.get("text")
                    if not isinstance(text, str) or not text.strip():
                        raise ValueError("paragraph.text must be a non-empty string")
                    p = body.add_paragraph() if body.text else body.paragraphs[0]
                    p.text = text
                    render_log.append(f"render:section[{idx}]:paragraph")
                    continue

                if block_type == "bullet_list":
                    bullets = block.get("bullets")
                    if not isinstance(bullets, list) or len(bullets) < 1:
                        raise ValueError("bullet_list.bullets must be a non-empty array")
                    for b in bullets:
                        if not isinstance(b, str):
                            raise ValueError("bullet_list.bullets items must be strings")
                        p = body.add_paragraph() if body.text else body.paragraphs[0]
                        p.text = b
                        p.level = 0
                    render_log.append(f"render:section[{idx}]:bullet_list")
                    continue

                if block_type == "table":
                    table = block.get("table")
                    if not isinstance(table, dict):
                        raise ValueError("table.table must be an object")
                    headers = table.get("headers")
                    rows = table.get("rows")
                    if not isinstance(headers, list) or not headers:
                        raise ValueError("table.headers must be a non-empty array")
                    if not isinstance(rows, list):
                        raise ValueError("table.rows must be an array")

                    # Minimal deterministic: render as plain text table
                    header_line = " | ".join([str(h) for h in headers])
                    p = body.add_paragraph() if body.text else body.paragraphs[0]
                    p.text = header_line
                    for r in rows:
                        if not isinstance(r, list):
                            raise ValueError("table.rows items must be arrays")
                        line = " | ".join([str(c) for c in r])
                        p = body.add_paragraph()
                        p.text = line
                    render_log.append(f"render:section[{idx}]:table")
                    continue

                if block_type == "citation":
                    cit = block.get("citation")
                    if not isinstance(cit, dict):
                        raise ValueError("citation.citation must be an object")
                    source = cit.get("source")
                    reference = cit.get("reference")
                    if not isinstance(source, str) or not source.strip():
                        raise ValueError("citation.source must be a non-empty string")
                    if reference is not None and not isinstance(reference, str):
                        raise ValueError("citation.reference must be a string")

                    note = s.notes_slide.notes_text_frame
                    note.text = f"citation: {source}\nreference: {reference or ''}".strip()
                    render_log.append(f"render:section[{idx}]:citation")

        prs.save(str(out_path))

        file_size = out_path.stat().st_size
        if file_size > max_file_size_mb * 1024 * 1024:
            out_path.unlink(missing_ok=True)
            raise ValueError("execution_constraints.max_file_size_mb exceeded")

        checksum = _sha256_file(out_path)
        return {
            "file_path": str(out_path),
            "slide_count": len(prs.slides),
            "checksum": checksum,
            "render_log": render_log,
        }
